"""
PYTHON DOCUMENT CHUNKER
Author: Gleb Khamin
Release Date: 23/10/2023
Version: Alpha 0.1
Description:
Python Library to chunk all documents in a directory structure.
Consists of two classes:
(1) DocumentChunker - chunks a single document and writes it to json
(2) DirectoryChunker - chunks all documents in a directory (and its subdirectories)
        using (1) above, and writes them all to a single /json_files directory
"""


import json
import os  # os.path provides access to the filesystem in Python
import sys  # for debugging sys.exit() is a quick way to exit the code
from pprint import pprint

# philopshy is: use well-supported python packages
# to decode the documents we want to chunk.
# Here are the packages we will use:
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook  # excel, note xlsx only (not xls - but could write a converter)

# the general design for the chunking is:
# (1) a class to chunk a single document and write it as json to file
# (2) a class to chunk a directory of documents, each time calling (1) above


# this is the class to chunk a SINGLE document
# suports multiple document types
class DocumentChunker:
    """
    This class chunks a single document into a dictionary/JSON type structure.
    Supported document types: .pptx, .pdf, .docx, .xlsx
    usage example:  dc = DocumentChunker("my_file.docx")
                    dc2 = DocumentChunker("anotherfile.xlsx")
                    chunks1 = dc.chunk_file()
                    chunks2 = dc2.chunk_file()



    """
    # This is the function that runs when you do: (for example) dc = DocumentChunker("my_file.docx")
    # Initialise the object, storing the filename to be decoded
    # note __init__ is pronounced "dunder init" because of the double underlines
    def __init__(self, filename):
        # in the __init__ we usually initialise the object's attributes
        # some attributes are set by the user, and some by the object.
        # When a user is setting up an attribute, it's normally passed to the
        # __init__ as an argument, e.g.
        # dc = DocumentChunker("my_file.docx")
        # passes "my_file.docx" to the __init__ function as the filename argument.
        # The self argument is always skipped when calling the __init__ function
        # using the classname. So we don't put:
        # dc = DocumentChunker(self, "my_file.docx")
        # we just put:
        # dc = DocumentChunker("my_file.docx")
        # BUT self always needs to be the first argument in the DEF line
        # of the __init__ function.
        # for example: filename is set by the user.
        self.filename = filename
        # but then the object, sets up a chunks attribute
        # that will be used by the methods later to store the chunks.
        # (initialise the dictionary that will be written to json)
        # we can initialise here in a useful way, because all document
        # types have the same top level JSON structure (e.g. "document_chunks":{}
        # document_type:"" etc)
        # we could've fille in document_type and original_filename in the line
        # below, but it's more readable to do it in two seperate lines below that.
        self.chunks = {"document_chunks":
                           {"document_type":"",
                            "original_filename":"",
                            "content": {}}}
        # in the __init__, as well as initilising the basic, we can add elements
        # to the basic initialised attributes.
        self.chunks["document_chunks"]['original_filename'] = os.path.basename(filename)
        self.chunks["document_chunks"]['document_type'] = filename.split(".")[-1]
        # the actual chunks will be placed in self.chunks["document_chunks"]["content"]

    def chunk_pptx_file(self):
        # pptx file has two levels of hierarchy:
        # slides and shapes
        pptx_dict = {'slides': []}
        # load in the pptx file
        pptx_file = Presentation(self.filename)
        for slide in pptx_file.slides:
            pptx_dict['slides'].append({'shapes': []})
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    if text not in ['‹#›', '', '\n']:
                        pptx_dict['slides'][-1]['shapes'].append({'text': text})
        return pptx_dict

    # can this be done more intelligently using the pdfreader package https://pypi.org/project/pdfreader/
    def chunk_pdf_file(self):
        reader = PdfReader(self.filename)
        pdf_dict = {'pages': []}
        for page in reader.pages:
            pdf_dict['pages'].append({'text': page.extract_text()})
        return pdf_dict


    # Called by parse_document() method to help with chunking docx files
    # function returns None if an error or not a heading
    # otherwise, returns the level of the heading
    def is_heading(self, paragraph, check_errors=True):
        if paragraph.style.name.startswith('Heading'):
            if check_errors:
                if len(paragraph.text) > 100:
                    print("Warning: heading is very long: ", paragraph.text)
                    return None
            return int(paragraph.style.name.split()[-1])


    # Called by chunk_docx_file() method - does the main work of chunking docx files
    def parse_document(self, document):
        #note , if a heading has only a table beneath, it may give empty text in the json for that heading's text content
        # or it could be a heading with a subheading directly below
        # which means it "contains no text for that heading" as such.
        structure = {'titled_subtitled_paragraphs':{}}
        joiner = "__"
        heading = ['']
        for paragraph in document.paragraphs:
            if paragraph.text.strip() == '':
                continue
            heading_level = self.is_heading(paragraph)
            if heading_level:   # python treats None as False, and numbers >0 as True
                # below here, the text must be a heading
                heading_level += 1  # this is just to make numbers easier to compare
                pt = paragraph.text
                # at the start, heading = ['content'], which is len = 1
                if heading_level > len(heading):  # is the only item in the heading list 'content'?
                    heading.append(pt.strip())
                elif heading_level == len(heading): # new same level heading
                    heading = heading[:-1] + [pt.strip()]
                else: # moving down a level
                    heading = heading[:heading_level - 1] + [pt.strip()]
                key = joiner.join(heading).replace("\t", "_").replace(" ", "_")
                if key.startswith("__"):
                    key = key[2:]
                if structure['titled_subtitled_paragraphs'].get(key) is None:
                    structure['titled_subtitled_paragraphs'][key] = ""
            else:   # just add content
                key = joiner.join(heading).replace("\t", "_").replace(" ", "_")
                if key.startswith("__"):
                    key = key[2:]
                if structure['titled_subtitled_paragraphs'].get(key) is None:
                    structure['titled_subtitled_paragraphs'][key] = ""
                structure['titled_subtitled_paragraphs'][key] += paragraph.text + "\n"
        tables = []
        if document.tables:
            for table in document.tables:
                table_json = {"columns": []}
                for col in table.columns:
                    col_json = {"col_text_list":[]}
                    for cell in col.cells:
                        if cell.text.strip() != '':
                            col_json["col_text_list"].append(cell.text)
                    if col_json["col_text_list"]:
                        table_json["columns"].append(col_json["col_text_list"])
                if table_json["columns"]:
                    tables.append(table_json)
            structure['tables'] = tables
        return structure

    def chunk_docx_file(self):
        # load in the docx file
        docx_file = Document(self.filename)
        return self.parse_document(docx_file)

    def chunk_xlsx_file(self):
        # note: a future change might be to infer the heading of a column
        # and make it the dictionary key for that column?
        # load in the xlsx file
        xlsx_file = load_workbook(self.filename)
        xlsx_chunks = {"sheets": []}
        sheets = xlsx_file.worksheets

        for sheet in sheets:
            sheet_json = {"sheet_name": sheet.title, "columns": []}
            for col in sheet.iter_cols(values_only=True):
                col_cells = []
                for cell in col:
                    # note: this will insert empty cells as empty strings
                    if cell is None:
                        col_cells.append("")
                        continue
                    col_cells.append(str(cell))
                sheet_json["columns"].append(col_cells)
            xlsx_chunks["sheets"].append(sheet_json)
        return xlsx_chunks

    def chunk_file(self):
        if self.filename.endswith(".pptx"):
            # each of these function / method calls below have two key elements:
            # 1. they are digging into the functions of the object and
            # 2. they are inserting the results into the chunks attribute of the object
            # remember self.chunks was initialised already in __init__
            self.chunks["document_chunks"]["content"] = self.chunk_pptx_file()
        elif self.filename.endswith(".pdf"):
            self.chunks["document_chunks"]["content"] = self.chunk_pdf_file()
        elif self.filename.endswith(".docx"):
            self.chunks["document_chunks"]["content"] = self.chunk_docx_file()
        elif self.filename.endswith(".xlsx"):
            self.chunks["document_chunks"]["content"] = self.chunk_xlsx_file()
        else:
            print("Unknown file type: " + self.filename)
        return self.chunks


# usage example:
# dtc = DirectoryTreeChunker("First batch of materials from Discy")
# dtc.chunk_all_files_in_directory_tree()
# (this writes all the json files to a directory called /json_files/ in the current directory)

class DirectoryTreeChunker:
    def __init__(self, directory):
        self.directory = directory
        self.filenames = self.get_all_docs_in_directory_tree(self.directory)

    # get all files in a directory and its subdirectories
    def get_all_files_directory_tree(self, directory):
        import os
        files = []
        for dirpath, dirnames, filenames in os.walk(directory):
            for filename in filenames:
                files.append(os.path.join(dirpath, filename))
        return files

    def is_doc_file(self, filename):
        doc_types = [".pptx", ".docx", ".xlsx", ".pdf"]
        for doc_type in doc_types:
            if filename.endswith(doc_type) and not os.path.basename(filename).startswith("~"):
                return True
        return False

    def get_all_docs_in_directory_tree(self, directory):
        files = self.get_all_files_directory_tree(directory)
        doc_files = []
        for file in files:
            if self.is_doc_file(file):
                doc_files.append(file)
        return doc_files

    @staticmethod
    def write_json_chunks(chunks_for_one_file, filename):
        fn = os.path.basename(filename)
        fn = "json_files/" + fn
        # json_files/...
        with open(fn + ".json", "w") as f:
            json.dump(chunks_for_one_file, f, indent=4)

    def chunk_all_files(self, write_files=True):
        if write_files:
            # create a directory to store the json files
            if not os.path.exists("json_files"):
                os.mkdir("json_files")
        chunks = {}
        for filename in self.filenames:
            print("Chunking file: ", filename)
            dc = DocumentChunker(filename)
            chunks[filename] = dc.chunk_file() # generates JSON for the file
            if write_files:
                self.write_json_chunks(chunks[filename], filename)
        return chunks


if __name__ == "__main__":
    """
    dc = DocumentChunker(".../First batch of materials from Discy/Doc examples/CX Example/01 Reference Files/CX Methodology.pdf")
    
    j = dc.chunk_file()
    print(j)
    """
    """p = DirectoryTreeChunker(".../First batch of materials from Discy")
    caf = p.chunk_all_files()
    with open('caf.json', 'w') as fp:
        json.dump(caf, fp, indent=4)"""
    """filename = ".../json_files/Process example.docx.json"
    with open(filename, "r") as f:
        j = json.load(f)
    #pprint(j)
    paras = j["document_chunks"]["content"]["titled_subtitled_paragraphs"]
    for para in paras.items():
        print(para)"""
    # testing just chunking a single file
    #filename = ".../First batch of materials from Discy/Doc examples/HBU and EBU Customer Journey Data_v1.xlsx"
    #filename = ".../First batch of materials from Discy/Doc examples/Process example/Process example.docx"
    #filename = ".../First batch of materials from Discy/Doc examples/North Sea - AI Delivery Review_ v0.3 (1).pptx"
    #filename = ".../First batch of materials from Discy/Doc examples/CX Example/01 Reference Files/CX Methodology.pdf"
    #dc = DocumentChunker(filename)
    #dc.chunk_file()
    #DirectoryTreeChunker.write_json_chunks(dc.chunks, filename)

    p = DirectoryTreeChunker("First batch of materials from Discy")
    caf = p.chunk_all_files()
    #with open('caf.json', 'w') as fp:
    #    json.dump(caf, fp, indent=4)
